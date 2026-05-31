#!/usr/bin/env python3

import re
import subprocess
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = REPO_ROOT / "presentation.pptx"
IMAGE_DIR = REPO_ROOT / "temp" / "slide-images"
HTML_PATH = REPO_ROOT / "index.html"
EXPECTED_SLIDES = 25
MAX_SIZE_BYTES = 200 * 1024 * 1024
VERIFY_ONLY = "--verify-only" in sys.argv


def clean_text(text):
    return re.sub(r"\s+", " ", text).strip()


def expected_titles():
    soup = BeautifulSoup(HTML_PATH.read_text(encoding="utf-8"), "lxml")
    titles = []
    for slide in soup.select(".slide"):
        title = slide.select_one(".slide-title")
        titles.append(clean_text(title.get_text(" ")) if title else "")
    return titles


if not VERIFY_ONLY:
    expected_images = [IMAGE_DIR / f"slide-{index:02d}.png" for index in range(1, EXPECTED_SLIDES + 1)]
    if not all(path.exists() for path in expected_images):
        subprocess.run(["node", "scripts/capture-slides.mjs"], cwd=REPO_ROOT, check=True)
    subprocess.run(["python3", "scripts/build-pptx.py"], cwd=REPO_ROOT, check=True)

assert PPTX_PATH.exists(), "presentation.pptx was not created"
assert PPTX_PATH.stat().st_size < MAX_SIZE_BYTES, "presentation.pptx is larger than 200MB"

presentation = Presentation(PPTX_PATH)
assert len(presentation.slides) == EXPECTED_SLIDES, (
    f"Expected {EXPECTED_SLIDES} slides, got {len(presentation.slides)}"
)

for index, slide in enumerate(presentation.slides, start=1):
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
    assert pictures, f"Slide {index} does not contain a background image"

titles = expected_titles()
for index, (slide, title) in enumerate(zip(presentation.slides, titles), start=1):
    notes_text = clean_text(slide.notes_slide.notes_text_frame.text)
    assert len(notes_text) > 10, f"Slide {index} notes appear to be empty"
    if title:
        assert title in notes_text, f"Slide {index} notes do not include title: {title}"

print("PPTX tests passed")
