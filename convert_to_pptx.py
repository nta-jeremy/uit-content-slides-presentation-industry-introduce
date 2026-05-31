#!/usr/bin/env python3
"""
Convert HTML presentation to PPTX format.
Extracts slide content from index.html and generates a PowerPoint file.
"""

import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).parent.resolve()
HTML_FILE = BASE_DIR / "index.html"
OUTPUT_PPTX = BASE_DIR / "presentation.pptx"
IMAGES_DIR = BASE_DIR / "images"

# ---------------------------------------------------------------------------
# Theme colors (dark theme matching the original)
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x0A, 0x0A, 0x12)          # Very dark navy/black
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)        # Neon cyan
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xC4)      # Teal
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED = RGBColor(0xA0, 0xA0, 0xB0)
TEXT_GOLD = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def clean_text(text):
    """Normalize whitespace in extracted text."""
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def add_bg_shape(slide, color=BG_COLOR):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape

def add_text_box(slide, left, top, width, height, text,
                 font_size=Pt(14), font_color=TEXT_WHITE,
                 bold=False, italic=False, align=PP_ALIGN.LEFT,
                 font_name="Calibri", word_wrap=True):
    """Add a text box with consistent formatting."""
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return tf

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=Pt(14), font_color=TEXT_WHITE,
                    bullet_color=ACCENT_CYAN, font_name="Calibri"):
    """Add a bullet list."""
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.name = font_name
    return tf

def add_accent_bar(slide, left, top, height, width=Inches(0.04), color=ACCENT_CYAN):
    """Add a thin vertical accent bar (like the original design)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(slide, soup_slide):
    """Slide 1 and 25: Title / Thank-you slide."""
    add_bg_shape(slide)

    category = soup_slide.find('div', class_='slide-category')
    title = soup_slide.find(['h1', 'h2', 'h3'], class_='slide-title')

    cat_text = clean_text(category.get_text()) if category else ""
    title_text = clean_text(title.get_text()) if title else ""

    # Accent bar on the left
    add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(3.5), color=ACCENT_CYAN)

    # Category
    if cat_text:
        add_text_box(slide, Inches(1.0), Inches(1.8), Inches(7.0), Inches(0.5),
                     cat_text, font_size=Pt(14), font_color=ACCENT_CYAN,
                     bold=True, align=PP_ALIGN.LEFT)

    # Title
    if title_text:
        add_text_box(slide, Inches(1.0), Inches(2.4), Inches(7.0), Inches(1.5),
                     title_text, font_size=Pt(40), font_color=TEXT_WHITE,
                     bold=True, align=PP_ALIGN.LEFT)

    # Extract paragraphs (team info, etc.)
    content_div = soup_slide.find('div', class_='slide-split') or soup_slide
    paragraphs = content_div.find_all('p')
    y_pos = Inches(4.2)
    for p in paragraphs:
        txt = clean_text(p.get_text())
        if txt and len(txt) > 2 and y_pos < Inches(6.5):
            color = TEXT_WHITE if p.find('strong') or p.find('b') else TEXT_MUTED
            sz = Pt(13) if p.find_parent('li') else Pt(14)
            add_text_box(slide, Inches(1.0), y_pos, Inches(7.0), Inches(0.4),
                         txt, font_size=sz, font_color=color)
            y_pos += Inches(0.35)

    # Image on the right side
    img = soup_slide.find('img')
    if img:
        src = img.get('src', '')
        img_path = BASE_DIR / src
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(8.5), Inches(1.8),
                                     width=Inches(4.0))

    # Footer
    footer = soup_slide.find('div', class_='slide-footer')
    if footer:
        footer_text = clean_text(footer.get_text())
        add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.3),
                     footer_text, font_size=Pt(10), font_color=TEXT_MUTED,
                     align=PP_ALIGN.LEFT)


def build_toc_slide(slide, soup_slide):
    """Slide 2: Table of Contents."""
    add_bg_shape(slide)

    category = soup_slide.find('div', class_='slide-category')
    title = soup_slide.find(['h1', 'h2', 'h3'], class_='slide-title')

    cat_text = clean_text(category.get_text()) if category else ""
    title_text = clean_text(title.get_text()) if title else ""

    add_accent_bar(slide, Inches(0.8), Inches(0.8), Inches(0.8), color=ACCENT_CYAN)

    if cat_text:
        add_text_box(slide, Inches(1.0), Inches(0.8), Inches(8.0), Inches(0.5),
                     cat_text, font_size=Pt(14), font_color=ACCENT_CYAN, bold=True)

    if title_text:
        add_text_box(slide, Inches(1.0), Inches(1.3), Inches(8.0), Inches(0.8),
                     title_text, font_size=Pt(36), font_color=TEXT_WHITE, bold=True)

    # Extract TOC items from grid
    grid_items = soup_slide.find_all('div', class_=re.compile('toc-item|grid-item'))
    if not grid_items:
        # Try finding numbered items
        grid_items = soup_slide.find_all('div', class_=re.compile('.*'))
        grid_items = [g for g in grid_items if re.search(r'\d+\.', clean_text(g.get_text()))]

    cols = 2
    rows = (len(grid_items) + cols - 1) // cols
    item_width = Inches(5.5)
    item_height = Inches(0.9)
    start_x = Inches(0.8)
    start_y = Inches(2.3)
    gap_x = Inches(0.4)
    gap_y = Inches(0.15)

    for i, item in enumerate(grid_items[:12]):
        col = i % cols
        row = i // cols
        x = start_x + col * (item_width + gap_x)
        y = start_y + row * (item_height + gap_y)

        txt = clean_text(item.get_text())
        if not txt:
            continue

        # Number extraction
        m = re.match(r'(\d+)\.\s*(.*)', txt)
        if m:
            num, rest = m.group(1), m.group(2)
            add_accent_bar(slide, x, y + Inches(0.05), Inches(0.04), Inches(0.5), color=ACCENT_TEAL)
            add_text_box(slide, x + Inches(0.1), y, item_width, Inches(0.4),
                         f"{num}. {rest}", font_size=Pt(16), font_color=TEXT_WHITE, bold=True)
        else:
            add_text_box(slide, x, y, item_width, Inches(0.4),
                         txt, font_size=Pt(15), font_color=TEXT_WHITE)


def build_content_slide(slide, soup_slide, slide_idx):
    """Generic content slide builder."""
    add_bg_shape(slide)

    category = soup_slide.find('div', class_='slide-category')
    title = soup_slide.find(['h1', 'h2', 'h3'], class_='slide-title')

    cat_text = clean_text(category.get_text()) if category else ""
    title_text = clean_text(title.get_text()) if title else ""

    # Header accent bar
    add_accent_bar(slide, Inches(0.5), Inches(0.5), Inches(0.6), color=ACCENT_CYAN)

    # Category label
    if cat_text:
        add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8.0), Inches(0.4),
                     cat_text, font_size=Pt(12), font_color=ACCENT_CYAN, bold=True)

    # Title
    title_y = Inches(0.9)
    if title_text:
        title_box = add_text_box(slide, Inches(0.7), title_y, Inches(8.5), Inches(0.9),
                                 title_text, font_size=Pt(32), font_color=TEXT_WHITE, bold=True)
        title_y = Inches(1.7)
    else:
        title_y = Inches(1.2)

    # Determine if there's an image
    img = soup_slide.find('img')
    has_image = False
    img_path = None
    if img:
        src = img.get('src', '')
        img_path = BASE_DIR / src
        if img_path.exists():
            has_image = True

    # Content area dimensions
    if has_image:
        content_width = Inches(6.5)
    else:
        content_width = Inches(11.5)

    content_left = Inches(0.7)

    # Extract list items
    lists = soup_slide.find_all(['ul', 'ol'])
    list_items = []
    for lst in lists:
        for li in lst.find_all('li', recursive=False):
            txt = clean_text(li.get_text())
            if txt and len(txt) > 2:
                list_items.append(txt)

    # Extract standalone paragraphs (not in lists, not category/title)
    standalone_paras = []
    all_paras = soup_slide.find_all('p')
    for p in all_paras:
        if p.find_parent(['li', 'div'], class_=re.compile('slide-category|slide-title|slide-footer')):
            continue
        txt = clean_text(p.get_text())
        if txt and len(txt) > 3 and txt not in [cat_text, title_text]:
            standalone_paras.append(txt)

    # Also extract h3/h4 headers
    headers = []
    for h in soup_slide.find_all(['h3', 'h4']):
        if 'slide-title' in (h.get('class') or []):
            continue
        txt = clean_text(h.get_text())
        if txt and len(txt) > 2:
            headers.append(txt)

    # Build content
    y_pos = title_y
    line_height = Inches(0.38)

    # Headers first
    for hdr in headers[:3]:
        if y_pos > Inches(6.5):
            break
        add_text_box(slide, content_left, y_pos, content_width, Inches(0.4),
                     hdr, font_size=Pt(16), font_color=ACCENT_TEAL, bold=True)
        y_pos += line_height

    # Bullet list
    if list_items:
        for item in list_items[:12]:
            if y_pos > Inches(6.7):
                break
            add_text_box(slide, content_left, y_pos, content_width, Inches(0.4),
                         f"• {item}", font_size=Pt(14), font_color=TEXT_WHITE)
            y_pos += line_height * 0.9

    # Standalone paragraphs
    for para in standalone_paras[:8]:
        if y_pos > Inches(6.7):
            break
        add_text_box(slide, content_left, y_pos, content_width, Inches(0.4),
                     para, font_size=Pt(13), font_color=TEXT_MUTED)
        y_pos += line_height * 0.9

    # Image on the right
    if has_image and img_path:
        try:
            slide.shapes.add_picture(str(img_path), Inches(7.5), Inches(1.5),
                                     width=Inches(5.0))
        except Exception as e:
            print(f"  Warning: Could not add image {img_path}: {e}")

    # Footer
    footer = soup_slide.find('div', class_='slide-footer')
    if footer:
        footer_text = clean_text(footer.get_text())
        add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.3),
                     footer_text, font_size=Pt(10), font_color=TEXT_MUTED, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Loading HTML...")
    with open(HTML_FILE, 'r', encoding='utf-8') as f:
        html = f.read()

    soup = BeautifulSoup(html, 'lxml')
    html_slides = soup.find_all('div', class_='slide')
    print(f"Found {len(html_slides)} slides")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for idx, html_slide in enumerate(html_slides):
        slide_num = idx + 1
        print(f"Processing slide {slide_num}...")

        slide = prs.slides.add_slide(blank_layout)

        if slide_num == 1 or slide_num == 25:
            build_title_slide(slide, html_slide)
        elif slide_num == 2:
            build_toc_slide(slide, html_slide)
        else:
            build_content_slide(slide, html_slide, slide_num)

    prs.save(OUTPUT_PPTX)
    print(f"\n✅ Saved to: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
