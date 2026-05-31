#!/usr/bin/env python3
"""Build a PowerPoint file from captured slide screenshots and HTML notes."""

from pathlib import Path
import re
import sys

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
HTML_PATH = BASE_DIR / "index.html"
IMAGE_DIR = BASE_DIR / "temp" / "slide-images"
OUTPUT_PATH = BASE_DIR / "presentation.pptx"
SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)


def clean_text(text):
    return re.sub(r"\s+", " ", text).strip()


def append_text(parts, text, seen):
    text = clean_text(text)
    if text and text not in seen:
        parts.append(text)
        seen.add(text)


def extract_slide_text(html_slide):
    parts = []
    seen = set()

    category = html_slide.find("div", class_="slide-category")
    title = html_slide.find(["h1", "h2", "h3"], class_="slide-title")
    if category:
        append_text(parts, f"[{category.get_text(' ')}]", seen)
    if title:
        append_text(parts, title.get_text(" "), seen)

    for selector in [
        "li",
        ".grid-item",
        ".feature-card",
        ".stat-card",
        ".timeline-item",
        ".quote-box",
    ]:
        for node in html_slide.select(selector):
            text = clean_text(node.get_text(" "))
            if text and text not in seen:
                parts.append(f"- {text}")
                seen.add(text)

    for paragraph in html_slide.find_all("p"):
        if paragraph.find_parent("li"):
            continue
        append_text(parts, paragraph.get_text(" "), seen)

    footer = html_slide.find("div", class_="slide-footer")
    if footer:
        append_text(parts, f"--- {footer.get_text(' ')} ---", seen)

    return "\n".join(parts)


def load_slides():
    if not HTML_PATH.exists():
        raise FileNotFoundError(f"Missing HTML source: {HTML_PATH}")

    soup = BeautifulSoup(HTML_PATH.read_text(encoding="utf-8"), "lxml")
    slides = soup.select(".slide")
    if not slides:
        raise ValueError("No .slide elements found in index.html")
    return slides


def add_notes(slide, text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    paragraph = notes_frame.paragraphs[0]
    paragraph.text = text
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(11)


def build_presentation():
    html_slides = load_slides()
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    blank_layout = presentation.slide_layouts[6]

    for index, html_slide in enumerate(html_slides, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        image_path = IMAGE_DIR / f"slide-{index:02d}.png"
        if not image_path.exists():
            raise FileNotFoundError(f"Missing slide image: {image_path}")

        slide.shapes.add_picture(
            str(image_path),
            Inches(0),
            Inches(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )
        add_notes(slide, extract_slide_text(html_slide))
        print(f"Slide {index:02d}: image + notes")

    presentation.save(OUTPUT_PATH)
    return len(html_slides)


def main():
    try:
        slide_count = build_presentation()
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    print(f"Built {OUTPUT_PATH} with {slide_count} slides")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
